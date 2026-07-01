#!/usr/bin/env python3
"""
HTML to PPTX Converter for hamada_presentation.html
Converts the HTML slide deck into an editable PowerPoint presentation.
"""

import re
from pathlib import Path
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Color scheme from HTML CSS
COLORS = {
    'blue': RGBColor(31, 78, 121),        # #1F4E79
    'blue_light': RGBColor(46, 117, 182),  # #2E75B6
    'blue_pale': RGBColor(222, 235, 247),  # #DEEBF7
    'accent': RGBColor(192, 0, 0),         # #C00000
    'gray': RGBColor(89, 89, 89),          # #595959
    'white': RGBColor(255, 255, 255),
    'text': RGBColor(34, 34, 34),          # #222
}

# Slide dimensions (16:9, 1280x720 in HTML -> standard PPTX 10"x5.625")
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)

# Standard margins
MARGIN_LEFT = Inches(0.5)
MARGIN_RIGHT = Inches(0.5)
MARGIN_TOP = Inches(0.4)
MARGIN_BOTTOM = Inches(0.45)


def extract_text_with_formatting(element):
    """Extract text and handle formatting tags (span, strong, em, etc.)"""
    if element is None:
        return ""
    
    if isinstance(element, str):
        return element.strip()
    
    result = []
    for child in element.children:
        if isinstance(child, str):
            result.append(child.strip())
        elif child.name in ['br']:
            result.append("\n")
        elif child.name in ['span', 'strong', 'em', 'b', 'i']:
            result.append(extract_text_with_formatting(child))
    
    return " ".join(result).strip()


def add_text_to_shape(shape, text, font_size=24, bold=False, color=None, align=PP_ALIGN.LEFT):
    """Add text to a shape with formatting."""
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.clear()
    
    if not text:
        return
    
    # Handle newlines
    for line in text.split('\n'):
        if not text_frame.paragraphs[0].text:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = line.strip()
        p.font.size = Pt(font_size)
        p.font.bold = bold
        if color:
            p.font.color.rgb = color
        p.alignment = align
    
    return text_frame


def process_slide(prs, slide_html, slide_num):
    """Process a single slide section and add to presentation."""
    # Create blank slide
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Set background color to white
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']
    
    # Check if it's a title slide
    is_title_slide = 'title-slide' in slide_html.get('class', [])
    
    if is_title_slide:
        process_title_slide(slide, slide_html, slide_num)
    else:
        process_content_slide(slide, slide_html, slide_num)


def process_title_slide(slide, slide_html, slide_num):
    """Process title/cover slide."""
    # Add top bar
    top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['blue']
    top_bar.line.color.rgb = COLORS['blue']
    
    left = MARGIN_LEFT
    top = Inches(0.8)
    width = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
    
    # Main title
    title_box = slide.shapes.add_textbox(left, top, width, Inches(2.0))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    
    title = slide_html.find('h1', class_='title')
    if title:
        title_text = extract_text_with_formatting(title)
        p.text = title_text
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = COLORS['blue']
        p.level = 0
    
    # Subtitle (Japanese)
    top += Inches(1.4)
    subtitle_box = slide.shapes.add_textbox(left, top, width, Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    
    subtitle = slide_html.find('h1', class_='subtitle-jp')
    if subtitle:
        subtitle_text = extract_text_with_formatting(subtitle)
        p.text = subtitle_text
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['gray']
        p.level = 0
    
    # Metadata
    top += Inches(1.0)
    meta_box = slide.shapes.add_textbox(left, top, width, Inches(1.2))
    meta_frame = meta_box.text_frame
    meta_frame.word_wrap = True
    
    meta = slide_html.find('div', class_='meta')
    if meta:
        for line in extract_text_with_formatting(meta).split('\n'):
            if meta_frame.paragraphs[0].text:
                p = meta_frame.add_paragraph()
            else:
                p = meta_frame.paragraphs[0]
            p.text = line.strip()
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['text']
    
    # Lab affiliation
    lab = slide_html.find('div', class_='lab')
    if lab:
        p = meta_frame.add_paragraph()
        p.text = extract_text_with_formatting(lab)
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['gray']
    
    # Date
    date = slide_html.find('div', class_='date')
    if date:
        p = meta_frame.add_paragraph()
        p.text = extract_text_with_formatting(date)
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['gray']


def process_content_slide(slide, slide_html, slide_num):
    """Process regular content slide."""
    # Add top bar
    top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['blue']
    top_bar.line.color.rgb = COLORS['blue']
    
    left = MARGIN_LEFT
    top = MARGIN_TOP
    width = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
    
    # Title/Header
    header = slide_html.find('h2', class_='header')
    if header:
        header_box = slide.shapes.add_textbox(left, top, width, Inches(0.6))
        header_frame = header_box.text_frame
        header_frame.word_wrap = True
        p = header_frame.paragraphs[0]
        p.text = extract_text_with_formatting(header)
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLORS['blue']
        
        top += Inches(0.65)
    
    # Content area
    content_box = slide.shapes.add_textbox(left, top, width, SLIDE_HEIGHT - top - MARGIN_BOTTOM)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    # Process all content elements except header and slide number
    for element in slide_html.find_all(True):
        if element.parent != slide_html:
            continue
        if element.name == 'h2':
            continue
        if element.get('class') and any(c in element.get('class', []) for c in ['slide-num', 'footer-cite', 'topbar']):
            continue
        
        # Paragraphs
        if element.name == 'p':
            text = extract_text_with_formatting(element)
            if text:
                if content_frame.paragraphs[0].text:
                    p = content_frame.add_paragraph()
                else:
                    p = content_frame.paragraphs[0]
                p.text = text
                p.font.size = Pt(22)
                p.font.color.rgb = COLORS['text']
                p.space_before = Pt(6)
                p.space_after = Pt(6)
        
        # Lists
        elif element.name == 'ol' or element.name == 'ul':
            for li in element.find_all('li', recursive=False):
                text = extract_text_with_formatting(li)
                if text:
                    if content_frame.paragraphs[0].text:
                        p = content_frame.add_paragraph()
                    else:
                        p = content_frame.paragraphs[0]
                    p.text = text
                    p.font.size = Pt(20)
                    p.font.color.rgb = COLORS['text']
                    p.level = 0
                    p.space_after = Pt(4)
        
        # Highlight boxes
        elif 'hlbox' in element.get('class', []):
            text = extract_text_with_formatting(element)
            if text:
                if content_frame.paragraphs[0].text:
                    p = content_frame.add_paragraph()
                else:
                    p = content_frame.paragraphs[0]
                p.text = text
                p.font.size = Pt(22)
                p.font.color.rgb = COLORS['text']
                p.space_before = Pt(8)
                p.space_after = Pt(8)
        
        # Figure boxes
        elif 'figbox' in element.get('class', []):
            text = extract_text_with_formatting(element)
            if text:
                if content_frame.paragraphs[0].text:
                    p = content_frame.add_paragraph()
                else:
                    p = content_frame.paragraphs[0]
                p.text = "[Figure] " + text
                p.font.size = Pt(18)
                p.font.italic = True
                p.font.color.rgb = COLORS['gray']
                p.space_before = Pt(8)
                p.space_after = Pt(8)
        
        # Tables
        elif element.name == 'table':
            add_table_to_slide(slide, element, left, top + Inches(2.5), width)


def add_table_to_slide(slide, table_html, left, top, width):
    """Add table to slide."""
    rows = table_html.find_all('tr')
    cols = len(rows[0].find_all(['th', 'td']))
    
    if not rows:
        return
    
    # Calculate dimensions
    table_height = Inches(min(2.5, len(rows) * 0.35))
    
    # Add table shape
    table_shape = slide.shapes.add_table(len(rows), cols, left, top, width, table_height).table
    
    # Fill table with data
    for row_idx, row_html in enumerate(rows):
        cells = row_html.find_all(['th', 'td'])
        for col_idx, cell_html in enumerate(cells):
            if col_idx < cols:
                cell = table_shape.cell(row_idx, col_idx)
                cell_text = extract_text_with_formatting(cell_html)
                
                # Set cell text
                text_frame = cell.text_frame
                text_frame.clear()
                if text_frame.paragraphs:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = cell_text
                p.font.size = Pt(14)
                
                # Header row formatting
                if row_html.name == 'tr' and row_html.find_parent('thead'):
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLORS['blue']
                    p.font.color.rgb = COLORS['white']
                    p.font.bold = True
                elif row_html.find('th'):
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLORS['blue']
                    p.font.color.rgb = COLORS['white']
                    p.font.bold = True


def convert_html_to_pptx(html_path, output_path):
    """Main conversion function."""
    # Read and parse HTML
    with open(html_path, 'r', encoding='utf-8') as f:
        html_content = f.read()
    
    soup = BeautifulSoup(html_content, 'html.parser')
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # Find all slides
    slides = soup.find_all('section', class_='slide')
    
    print(f"Found {len(slides)} slides")
    
    # Process each slide
    for slide_num, slide_html in enumerate(slides, 1):
        print(f"Processing slide {slide_num}/{len(slides)}...", end=' ')
        try:
            process_slide(prs, slide_html, slide_num)
            print("✓")
        except Exception as e:
            print(f"✗ Error: {e}")
    
    # Save presentation
    prs.save(output_path)
    print(f"\n✓ Presentation saved to: {output_path}")


if __name__ == '__main__':
    html_file = Path(__file__).parent / 'hamada_presentation.html'
    output_file = Path(__file__).parent / 'hamada_presentation.pptx'
    
    if not html_file.exists():
        print(f"Error: HTML file not found: {html_file}")
        exit(1)
    
    convert_html_to_pptx(str(html_file), str(output_file))
